import os
import pdfplumber
import docx
import pptx

def extract_text_from_file(file_obj, filename=""):
    ext = os.path.splitext(filename.lower())[1]
    
    text = ""
    try:
        if ext == ".docx":
            doc = docx.Document(file_obj)
            for para in doc.paragraphs:
                if para.text:
                    text += para.text + "\n"
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text += cell.text + " "
                    text += "\n"
        elif ext == ".pptx":
            prs = pptx.Presentation(file_obj)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        else:
            # Fallback to pdfplumber for .pdf or unknown types
            with pdfplumber.open(file_obj) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
    except Exception as e:
        # If all else fails, attempt to read as raw text
        try:
            file_obj.seek(0)
            raw_data = file_obj.read()
            text = raw_data.decode('utf-8')
        except Exception:
            raise ValueError(f"Failed to extract text from {filename or 'file'}: {str(e)}")
            
    return text
